"""
SmartRecruit AI Platform — Comprehensive PPTX Presentation Generator
Produces a 22-slide professional presentation.
Run: python generate_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG = lambda name: os.path.join(BASE_DIR, name)

# ── THEME COLORS ──────────────────────────
DARK_BG  = RGBColor(0x0A, 0x0B, 0x10)   # #0a0b10
ACCENT   = RGBColor(0x00, 0xD2, 0xFF)   # #00d2ff (cyan)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
INDIGO   = RGBColor(0x63, 0x66, 0xF1)   # #6366f1
GREEN    = RGBColor(0x10, 0xB9, 0x81)   # emerald
ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)   # amber

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide(bg_color=DARK_BG):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_image(slide, path, left, top, width=None, height=None):
    full = IMG(path) if not os.path.isabs(path) else path
    if os.path.exists(full):
        if width and height:
            slide.shapes.add_picture(full, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(full, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(full, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(full, Inches(left), Inches(top))

def divider(slide, color=ACCENT, top=1.4, alpha=False):
    add_rect(slide, 0.6, top, 3.0, 0.04, color)

def slide_header(slide, number, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.3, RGBColor(0x0F, 0x17, 0x28))
    add_text(slide, number, 0.6, 0.15, 11, 0.5, size=11, color=ACCENT, bold=True)
    add_text(slide, subtitle, 0.6, 0.55, 11, 0.6, size=22, bold=True, color=WHITE)
    add_rect(slide, 0, 1.28, 13.33, 0.03, ACCENT)

# ════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ════════════════════════════════════════════
s1 = add_slide()
add_rect(s1, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(s1, 0, 0, 0.35, 7.5, ACCENT)
add_rect(s1, 0.35, 0, 12.98, 0.04, INDIGO)

add_text(s1, "SMARTRECRUIT", 1.0, 1.5, 11, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "AI RECRUITMENT PLATFORM", 1.0, 2.7, 11, 0.7, size=24, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)
add_rect(s1, 1.0, 3.45, 6.0, 0.05, ACCENT)
add_text(s1, "A Next-Generation Talent Acquisition Ecosystem\nPowered by Google Gemini · Multi-lingual RAG · Real-time AI Interviews",
         1.0, 3.6, 11, 1.0, size=13, color=MUTED, align=PP_ALIGN.LEFT)
add_text(s1, "Post Graduate Dept. of Computer Science & Technology\nSardar Patel University · 2024–2026  |  IR Info Tech Pvt. Ltd.",
         1.0, 5.6, 11, 0.8, size=11, color=MUTED, italic=True, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT
# ════════════════════════════════════════════
s2 = add_slide()
slide_header(s2, "01", "The Problem We Solve")
divider(s2, ACCENT, top=1.35)

problems = [
    ("⏳", "Time-to-Hire", "Average 42 days globally. 70% spent manually screening resumes."),
    ("⚖️", "Hiring Bias", "Unconscious bias in manual screening leads to diverse-talent loss."),
    ("📉", "Poor Scale", "Hundreds of applicants per role; HR teams are overwhelmed."),
    ("🤖", "No AI in HR", "Most HR tools lack real AI — they use keyword matching, not intelligence."),
]
for i, (icon, title, desc) in enumerate(problems):
    x = 0.6 + (i % 2) * 6.4
    y = 1.6 + (i // 2) * 2.5
    add_rect(s2, x, y, 5.8, 2.0, RGBColor(0x14, 0x19, 0x28))
    add_text(s2, icon + " " + title, x+0.2, y+0.15, 5.4, 0.5, size=14, bold=True, color=ACCENT)
    add_text(s2, desc, x+0.2, y+0.65, 5.4, 1.1, size=11, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 3 — SOLUTION OVERVIEW
# ════════════════════════════════════════════
s3 = add_slide()
slide_header(s3, "02", "SmartRecruit: The Complete Solution")
divider(s3)

features = [
    ("🔍", "AI Resume Parser", "Multi-lingual RAG + Semantic Embeddings"),
    ("📝", "Auto Assessments", "AI-generated Aptitude + Practical MCQs"),
    ("🎙️", "AI Interview (R3)", "Voice, Code, Sentiment, Face Proctoring"),
    ("🌿", "Botanist HR (R4)", "Behavioral STAR-method AI evaluation"),
    ("📊", "Analytics", "Hiring funnel, retention prediction, bias report"),
    ("🎁", "Offer & Onboard", "PDF generation, salary negotiation AI, roadmap"),
]
for i, (icon, title, desc) in enumerate(features):
    x = 0.6 + (i % 3) * 4.2
    y = 1.6 + (i // 3) * 2.5
    add_rect(s3, x, y, 3.9, 2.0, RGBColor(0x14, 0x19, 0x28))
    add_rect(s3, x, y, 3.9, 0.06, ACCENT)
    add_text(s3, icon, x+0.2, y+0.2, 0.5, 0.5, size=20, color=WHITE)
    add_text(s3, title, x+0.8, y+0.22, 2.9, 0.45, size=13, bold=True, color=WHITE)
    add_text(s3, desc, x+0.2, y+0.85, 3.5, 0.9, size=10, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 4 — 4-ROUND PIPELINE
# ════════════════════════════════════════════
s4 = add_slide()
slide_header(s4, "03", "The 4-Round AI Hiring Pipeline")
divider(s4)

rounds = [
    (ACCENT,  "R1", "Aptitude Test", "AI-generated MCQs\nTimed | Anti-cheat\nPass: 70%"),
    (INDIGO,  "R2", "Practical Test", "Domain-specific MCQs\nCode submission\nPass: 70%"),
    (GREEN,   "R3", "AI Interview", "Voice + Code + Camera\nSentiment analysis\nPass: 75%"),
    (ORANGE,  "R4", "HR Interview", "Botanist STAR engine\nMulti-lingual voice\nRecruiter approved"),
]
for i, (color, label, title, desc) in enumerate(rounds):
    x = 0.6 + i * 3.1
    add_rect(s4, x, 1.6, 2.8, 4.6, RGBColor(0x14, 0x19, 0x28))
    add_rect(s4, x, 1.6, 2.8, 0.08, color)
    add_text(s4, label, x + 0.9, 1.75, 1.0, 0.7, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s4, title, x + 0.1, 2.5, 2.6, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, desc, x + 0.15, 3.1, 2.5, 2.0, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s4, "→", x + 2.8, 3.5, 0.3, 0.5, size=20, bold=True, color=ACCENT)

# ════════════════════════════════════════════
#  SLIDE 5 — TECHNOLOGY STACK
# ════════════════════════════════════════════
s5 = add_slide()
slide_header(s5, "04", "Technology Stack")
divider(s5)

layers = [
    ("Backend", "Python 3.11 · Django 4.x · Django REST Framework · JWT"),
    ("AI/ML Engines", "Google Gemini 1.5 Flash · Groq (LLaMA3) · HuggingFace (Mistral-7B)"),
    ("NLP/Embeddings", "text-embedding-004 · SentenceTransformers · langdetect · spaCy"),
    ("Real-time", "Django Channels · Redis · WebSockets"),
    ("Frontend", "Bootstrap 5 · Vanilla JS · Chart.js · Monaco Editor"),
    ("Authentication", "django-allauth · Azure AD · Okta SAML · bcrypt"),
    ("Database", "SQLite (dev) · PostgreSQL (prod) · Django ORM"),
    ("Deployment", "WhiteNoise · Gunicorn · Nginx · PWA Service Worker"),
]
for i, (layer, tech) in enumerate(layers):
    x = 0.6 + (i % 2) * 6.4
    y = 1.5 + (i // 2) * 1.4
    add_rect(s5, x, y, 5.9, 1.15, RGBColor(0x14, 0x19, 0x28))
    add_text(s5, layer, x+0.2, y+0.1, 5.5, 0.4, size=11, bold=True, color=ACCENT)
    add_text(s5, tech, x+0.2, y+0.52, 5.5, 0.5, size=10, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 6 — AI ENGINE ARCHITECTURE
# ════════════════════════════════════════════
s6 = add_slide()
slide_header(s6, "05", "AI Engine: 3-Tier Graceful Fallback")
divider(s6)
add_text(s6, "Zero-downtime: If Gemini hits quota, traffic auto-routes to Groq → Hugging Face", 0.6, 1.5, 12, 0.4, size=12, color=MUTED)
engines = [
    (ACCENT, "PRIMARY", "Google Gemini 1.5 Flash", "High-context AI\nResume parsing, JD generation,\nInterview evaluation, Architecture scoring"),
    (ORANGE, "FALLBACK 1", "Groq API (LLaMA3-8B)", "Ultra-low latency\nVoice assistant, Chat,\ntyping analysis"),
    (INDIGO, "FALLBACK 2", "HuggingFace Mistral-7B", "Ultimate offline backup\nAll tasks via\nHF Inference API"),
]
for i, (color, label, name, desc) in enumerate(engines):
    x = 0.6 + i * 4.2
    add_rect(s6, x, 2.1, 3.9, 4.2, RGBColor(0x14, 0x19, 0x28))
    add_rect(s6, x, 2.1, 3.9, 0.07, color)
    add_text(s6, label, x+0.15, 2.22, 3.6, 0.4, size=9, bold=True, color=color)
    add_text(s6, name, x+0.15, 2.65, 3.6, 0.5, size=13, bold=True, color=WHITE)
    add_text(s6, desc, x+0.15, 3.25, 3.6, 2.0, size=10, color=MUTED)
    if i < 2:
        add_text(s6, "↓ fallback", x+3.9, 3.8, 0.9, 0.5, size=9, color=ACCENT)

# ════════════════════════════════════════════
#  SLIDE 7 — USE CASE DIAGRAM
# ════════════════════════════════════════════
s7 = add_slide()
slide_header(s7, "06", "Use Case Diagram")
divider(s7)
add_text(s7, "33 use cases across 3 actors: Candidate, Recruiter, Admin", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s7, "use_case.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 8 — DFD LEVEL 0
# ════════════════════════════════════════════
s8 = add_slide()
slide_header(s8, "07", "Data Flow Diagram — Level 0 (Context)")
divider(s8)
add_text(s8, "System boundary with all 6 external entities", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s8, "dfd_level_0.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 9 — DFD LEVEL 1
# ════════════════════════════════════════════
s9 = add_slide()
slide_header(s9, "08", "Data Flow Diagram — Level 1 (Processes)")
divider(s9)
add_text(s9, "7 major sub-processes and 7 data stores", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s9, "dfd_level_1.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 10 — DFD LEVEL 2
# ════════════════════════════════════════════
s10 = add_slide()
slide_header(s10, "09", "DFD Level 2 — Resume Intelligence Deep-Dive")
divider(s10)
add_text(s10, "From PDF upload → multilingual detect → embed → score → shortlist/reject", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s10, "dfd_level_2.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 11 — ACTIVITY DIAGRAM
# ════════════════════════════════════════════
s11 = add_slide()
slide_header(s11, "10", "Activity Diagram — Full Hiring Pipeline")
divider(s11)
add_text(s11, "4-round elimination with all decision gates and email triggers", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s11, "activity_flow.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 12 — CLASS DIAGRAM
# ════════════════════════════════════════════
s12 = add_slide()
slide_header(s12, "11", "Class Diagram — Core Data Models")
divider(s12)
add_text(s12, "9 core classes with attributes, methods, and multiplicities", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s12, "class_diagram.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 13 — SEQUENCE DIAGRAM
# ════════════════════════════════════════════
s13 = add_slide()
slide_header(s13, "12", "Sequence Diagram — End-to-End Hiring Flow")
divider(s13)
add_text(s13, "Apply → Parse → Assess → Interview → Offer → Onboard — with async email & AI calls", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s13, "sequence_diagram.png", left=1.5, top=1.9, height=5.2)

# ════════════════════════════════════════════
#  SLIDE 14 — DATABASE DESIGN
# ════════════════════════════════════════════
s14 = add_slide()
slide_header(s14, "13", "Database Design: 35+ Models")
divider(s14)
models = [
    ("User", "core", "Custom auth model with role flags"),
    ("JobPosting", "jobs", "Full JD with round config"),
    ("Candidate", "jobs", "Profile + AI-enriched fields"),
    ("Application", "jobs", "22-status pipeline state machine"),
    ("Assessment", "jobs", "R1/R2 results + JSON details"),
    ("Interview", "jobs", "AI session data + code + scores"),
    ("Offer", "jobs", "Salary, joining date, PDF link"),
    ("SentimentLog", "jobs", "Per-frame emotion, proctoring flags"),
    ("Notification", "jobs", "In-app bell notifications"),
    ("NegotiationSession", "jobs", "AI salary negotiation chat"),
    ("CodingChallenge", "jobs", "XP, test cases, multi-language"),
    ("OnboardingRoadmap", "jobs", "30-60-90 day AI plan"),
]
for i, (model, app, desc) in enumerate(models):
    x = 0.6 + (i % 3) * 4.2
    y = 1.55 + (i // 3) * 1.55
    add_rect(s14, x, y, 3.9, 1.3, RGBColor(0x14, 0x19, 0x28))
    add_text(s14, model, x+0.2, y+0.1, 2.5, 0.4, size=12, bold=True, color=ACCENT)
    add_text(s14, f"[{app}]", x+2.6, y+0.12, 1.1, 0.35, size=9, color=MUTED)
    add_text(s14, desc, x+0.2, y+0.6, 3.5, 0.55, size=9, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 15 — SECURITY
# ════════════════════════════════════════════
s15 = add_slide()
slide_header(s15, "14", "Security Architecture")
divider(s15)

sec = [
    ("🔐", "JWT Authentication", "15-min expiry access tokens + refresh"),
    ("🛡️", "CSRF Protection", "Built-in Django middleware on all forms"),
    ("💉", "SQL Injection", "ORM only — zero raw SQL strings"),
    ("🔒", "HTTPS Enforced", "HSTS max-age 1 year + SSL redirect"),
    ("🎭", "Blind Hiring", "Toggle masks candidate name/photo"),
    ("📸", "Proctoring", "Face detection, tab-switch, flag system"),
    ("⛓️", "Blockchain", "SHA-256 credential verification"),
    ("📁", "File Safety", "MIME + extension whitelist, 5 MB cap"),
]
for i, (icon, title, desc) in enumerate(sec):
    x = 0.6 + (i % 2) * 6.4
    y = 1.6 + (i // 2) * 1.4
    add_rect(s15, x, y, 5.9, 1.15, RGBColor(0x14, 0x19, 0x28))
    add_text(s15, icon + " " + title, x+0.2, y+0.1, 5.5, 0.45, size=12, bold=True, color=GREEN)
    add_text(s15, desc, x+0.2, y+0.6, 5.5, 0.45, size=10, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 16 — LANDING PAGE SCREENSHOT
# ════════════════════════════════════════════
s16 = add_slide()
slide_header(s16, "15", "System UI: Landing Page")
divider(s16)
add_text(s16, "Premium glassmorphism design with Dark / Light theme switching", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s16, "screenshot_landing.png", left=0.5, top=1.85, height=5.4)

# ════════════════════════════════════════════
#  SLIDE 17 — JOB MATRIX SCREENSHOT
# ════════════════════════════════════════════
s17 = add_slide()
slide_header(s17, "16", "System UI: Neural Job Matrix")
divider(s17)
add_text(s17, "AI match scores, advanced filters, pagination — candidate-facing job board", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s17, "screenshot_jobs.png", left=0.5, top=1.85, height=5.4)

# ════════════════════════════════════════════
#  SLIDE 18 — REGISTRATION SCREENSHOT
# ════════════════════════════════════════════
s18 = add_slide()
slide_header(s18, "17", "System UI: Candidate Registration")
divider(s18)
add_text(s18, "Glassmorphism form with animated floating labels, SSO, and PWA install prompt", 0.6, 1.5, 12, 0.35, size=12, color=MUTED)
add_image(s18, "screenshot_register.png", left=0.5, top=1.85, height=5.4)

# ════════════════════════════════════════════
#  SLIDE 19 — KEY ACHIEVEMENTS
# ════════════════════════════════════════════
s19 = add_slide()
slide_header(s19, "18", "Key Achievements")
divider(s19)
items = [
    ("✅", "0 Django Errors", "python manage.py check — 0 issues silenced"),
    ("✅", "35+ Data Models", "Complete hiring pipeline state machine"),
    ("✅", "3-Engine AI", "Gemini → Groq → HuggingFace graceful fallback"),
    ("✅", "Multi-lingual", "30+ languages via multilingual embeddings"),
    ("✅", "Real-time", "WebSocket proctoring, live interview, chat"),
    ("✅", "PWA Ready", "Offline cache, install prompt, service worker"),
    ("✅", "SSO Enterprise", "Azure AD + Okta SAML + JWT"),
    ("✅", "Blockchain", "Credential tamper-proof verification"),
    ("✅", "Dark + Light", "Full dual-theme CSS system (4000+ lines)"),
    ("✅", "Auto Emails", "Branded HTML emails via async threading"),
]
for i, (icon, title, desc) in enumerate(items):
    x = 0.6 + (i % 2) * 6.4
    y = 1.5 + (i // 2) * 1.15
    add_text(s19, icon + " " + title + " — " + desc, x, y, 6.0, 0.9, size=11, color=WHITE)

# ════════════════════════════════════════════
#  SLIDE 20 — FUTURE SCOPE
# ════════════════════════════════════════════
s20 = add_slide()
slide_header(s20, "19", "Future Scope")
divider(s20)
future = [
    "LinkedIn / Naukri / Indeed auto job syndication",
    "3D face liveness detection (anti-spoofing proctoring)",
    "RL-based question difficulty adaptation per candidate",
    "Multi-tenant SaaS white-labelling for enterprises",
    "Predictive HR analytics with regression ML models",
    "GDPR right-to-erasure automation and data anonymization",
]
for i, f in enumerate(future):
    add_text(s20, f"▸  {f}", 0.6, 1.7 + i * 0.85, 12, 0.7, size=13, color=MUTED)

# ════════════════════════════════════════════
#  SLIDE 21 — CONCLUSION
# ════════════════════════════════════════════
s21 = add_slide()
slide_header(s21, "20", "Conclusion")
divider(s21)
add_text(
    s21,
    "SmartRecruit successfully demonstrates that Artificial Intelligence can be applied "
    "end-to-end across the complete talent acquisition lifecycle — from resume parsing to "
    "onboarding roadmap generation.\n\n"
    "The platform reduces time-to-shortlist from days to seconds, eliminates human bias from "
    "the screening process, and delivers a world-class, data-driven experience to both "
    "candidates and recruiters.\n\n"
    "With a three-tier zero-downtime AI engine, 4-round automated evaluation pipeline, "
    "and enterprise-grade security — SmartRecruit is production-ready.",
    0.6, 1.6, 12.1, 5.5, size=14, color=MUTED
)

# ════════════════════════════════════════════
#  SLIDE 22 — THANK YOU
# ════════════════════════════════════════════
s22 = add_slide()
add_rect(s22, 0, 0, 0.35, 7.5, ACCENT)
add_rect(s22, 0.35, 0, 12.98, 0.04, INDIGO)
add_text(s22, "Thank You", 1.0, 2.0, 11, 1.5, size=52, bold=True, color=WHITE)
add_text(s22, "SmartRecruit AI Recruitment Platform", 1.0, 3.6, 11, 0.6, size=18, color=ACCENT)
add_rect(s22, 1.0, 4.3, 5.0, 0.05, ACCENT)
add_text(s22, "Questions & Discussion", 1.0, 4.5, 11, 0.5, size=14, color=MUTED, italic=True)

# ════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════
out = os.path.join(BASE_DIR, "SmartRecruit_Final_Presentation_Compliant.pptx")
prs.save(out)
print(f"[OK] PPTX saved -> {out}")
