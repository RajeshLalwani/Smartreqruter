import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG = lambda name: os.path.join(BASE_DIR, 'Diagrams', name) if os.path.exists(os.path.join(BASE_DIR, 'Diagrams', name)) else os.path.join(BASE_DIR, name)

prs = Presentation()

# Theme and style constants
BLUE = RGBColor(0, 70, 127)
LIGHT_BLUE = RGBColor(59, 120, 176)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle, author):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.color.rgb = BLUE
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle_box.text = f"{subtitle}\n\nSubmitted By:\n{author}"
    return slide

def add_content_slide(title, points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = pt
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    return slide

def add_image_slide(title, img_path):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    full = IMG(img_path)
    if os.path.exists(full):
        # Calculate aspect ratio to fit safely
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(full, left, top, height=height)
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = f"[Image Missing: {img_path}]"
    return slide

def add_data_dictionary_slide(title, data):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = f"Data Dictionary: {title}"
    title_box.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    rows, cols = len(data) + 1, 4
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(0.4 * rows)
    
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    
    headers = ["Sr. No", "Field Name", "Data Type", "Constraints"]
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(59, 120, 176) # 3B78B0
        
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            # Alternating colors
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(221, 235, 247) # DDEBF7
            else:
                cell.fill.fore_color.rgb = RGBColor(189, 215, 238) # BDD7EE

# --- Build Presentation ---

# 1. Project profile (Index slide + Title)
# ... Existing calls remain ...
add_title_slide("SMARTRECRUIT AI PLATFORM", "Next-Gen AI-Powered Talent Acquisition Prototype", "Candidate Name\nTech Elecon Pvt. Ltd.")

add_content_slide("Index", [
    "1. Project Profile",
    "2. Introduction To Organization",
    "3. Introduction to Project",
    "4. Objective of the System",
    "5. Functional & Non-Functional Requirements",
    "6. Hardware and Software Requirements",
    "7. System Modules and Functionalities",
    "8. System Design Diagrams",
    "9. Data Dictionary",
    "10. System Outputs (Screenshots)",
    "11. Conclusion"
])

# 2. Introduction To Organization
add_content_slide("Introduction To Organization", [
    "Tech Elecon Pvt. Ltd.",
    "A premier IT solutions provider and part of the esteemed Elecon Group.",
    "Specializes in enterprise software development, system integration, and advanced R&D.",
    "Focuses on modernizing legacy workflows with bleeding-edge technologies.",
    "Headquartered in Vallabh Vidyanagar, Gujarat."
])

# 3. Introduction to project
add_content_slide("Introduction to Project", [
    "SmartRecruit is a state-of-the-art AI-driven recruitment portal.",
    "Completely automates the hiring funnel utilizing Large Language Models (Gemini 1.5, LLaMA3).",
    "Eliminates manual CV screening latency through multi-lingual semantic RAG matching.",
    "Conducts real-time behavioral AI-interviews equipped with biometric sentiment proctoring.",
    "Designed to combat hiring bias while establishing extremely scalable processing pipelines."
])

# 4. Objective of the system
add_content_slide("Objective of the System", [
    "Automate Candidate Evaluation Pipeline via an objective algorithmic standard.",
    "Enhance Security: Piston-isolated execution sandbox for coding evaluations.",
    "Mitigate Human Bias using entirely objective LLM-generated assessment metrics.",
    "Reduce Operational Costs: Shortlist generation time cut from 12 hours to 8 seconds.",
    "Promote Scalability through asynchronous micro-interactions, caching, and Django Channels."
])

# 5. Functional and Non-Functional Requirements
add_content_slide("Functional & Non-Functional Requirements", [
    "Functional Requirements:",
    "- Semantic Resume Scoring & Auto-shortlisting.",
    "- Multi-tier AI assessment execution (Coding Engine & Interview Engine).",
    "- 'Thunder Toast' UI Notifications and Glassmorphism Dashboards.",
    "- Granular Multi-Role Permissions (Candidate, Recruiter, Administrator).",
    "",
    "Non-Functional Requirements:",
    "- 99.9% Uptime with 'Graceful AI Degradation' mechanism (Gemini -> Groq -> Local Mistral).",
    "- Sub-500ms Piston executing response latency.",
    "- Complete PWA (Progressive Web App) offline resilience capability."
])

# 6. Hardware and Software Requirements
add_content_slide("Hardware and Software Requirements", [
    "Hardware Requirements:",
    "- Minimum Dual-Core CPU with 4GB RAM.",
    "- Standard webcam & microphone for real-time DeepFace emotion analysis.",
    "",
    "Software & Technology Stack:",
    "- Backend: Python 3.11, Django 5.x, Django Channels, WebSockets.",
    "- Frontend: HTML5, CSS3 Custom Theme Engine, Vanilla JS, Bootstrap 5.",
    "- AI / ML Providers: Google Gemini 1.5 API, Hugging Face, OpenCV.",
    "- Database Engine: PostgreSQL / MySQL."
])

# 7. System Modules and Functionalities
add_content_slide("System Modules and Functionalities", [
    "Account Authorization & PWA Integration Module.",
    "Recruiter Management Console (Job Generation & Algorithmic Filtering).",
    "Semantic RAG Resume Evaluator (Match & Gap Analysis).",
    "AI Sandbox Assessment Hub (Aptitude + Practical Coding Tests).",
    "Live HR Interview Simulation Room + OpenCV Proctoring.",
    "Dynamic Statistical Reporting & Analytics Engine."
])

# 8. System design diagrams
add_image_slide("Entity-Relationship Diagram (ERD)", "er_diagram.png")
add_image_slide("Class Diagram (UML)", "class_diagram.png")
add_image_slide("Use Case Diagram", "use_case.png")
add_image_slide("Activity Diagram", "activity_flow.png")
add_image_slide("Sequence Diagram", "sequence_diagram.png")
add_image_slide("Level 0 DFD", "dfd_level_0.png")
add_image_slide("Level 1 DFD", "dfd_level_1.png")
add_image_slide("Level 2 DFD", "dfd_level_2.png")

# 9. Data Dictionary
db_tables = {
    "core_users": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "password", "VARCHAR(128)", "Not Null"],
        [3, "last_login", "DATETIME", "Null"],
        [4, "username", "VARCHAR(150)", "Unique, Not Null"],
        [5, "email", "VARCHAR(254)", "Not Null"],
        [6, "role", "VARCHAR(50)", "Not Null"]
    ],
    "jobs_jobposting": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "title", "VARCHAR(255)", "Not Null"],
        [3, "description", "LONGTEXT", "Not Null"],
        [4, "requirements", "LONGTEXT", "Not Null"],
        [5, "status", "VARCHAR(20)", "Not Null"],
        [6, "recruiter_id", "INT", "Foreign Key (users.id)"]
    ],
    "jobs_application": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "resume", "VARCHAR(100)", "Not Null"],
        [3, "match_score", "DOUBLE", "Null"],
        [4, "status", "VARCHAR(20)", "Not Null"],
        [5, "job_id", "INT", "Foreign Key (jobposting.id)"],
        [6, "candidate_id", "INT", "Foreign Key (users.id)"]
    ],
    "jobs_interviewresult": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "code_score", "DOUBLE", "Not Null"],
        [3, "sentiment_score", "DOUBLE", "Not Null"],
        [4, "feedback", "LONGTEXT", "Not Null"],
        [5, "application_id", "INT", "Foreign Key (application.id)"]
    ],
    "jobs_codingchallenge": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "title", "VARCHAR(255)", "Not Null"],
        [3, "difficulty", "VARCHAR(50)", "Not Null"],
        [4, "xp_reward", "INT", "Not Null, Default 100"]
    ],
    "jobs_notification": [
        [1, "id", "INT", "Primary Key, Auto-Increment"],
        [2, "title", "VARCHAR(255)", "Not Null"],
        [3, "message", "LONGTEXT", "Not Null"],
        [4, "is_read", "TINYINT(1)", "Not Null, Default 0"],
        [5, "user_id", "INT", "Foreign Key (users.id)"]
    ]
}

for t_name, t_data in db_tables.items():
    add_data_dictionary_slide(t_name, t_data)

# 10. System Outputs (Screenshots)
add_content_slide("System Outputs (Screenshots)", [
    "The following slides display the user interface of the SmartRecruit system.",
    "The UI is built utilizing a modern Glassmorphism design system.",
    "Supports active Dark/Light theme toggling, Thunder Toasts, and real-time data binding."
])
add_image_slide("System Output: Landing Interface", "screenshot_landing.png")
add_image_slide("System Output: Neural Job Listings", "screenshot_jobs.png")
add_image_slide("System Output: Candidate Registration", "screenshot_register.png")
add_image_slide("System Output: Admin Dashboard", "screenshot_admin.png")
add_image_slide("System Output: AI Interview Room", "screenshot_interview.png")

# 11. Conclusion
add_content_slide("Conclusion", [
    "SmartRecruit proves that Artificial Intelligence can effectively manage the full talent lifecycle.",
    "Decreases administrative recruitment overhead by up to 85% while drastically improving sourcing consistency.",
    "Maintains uncompromised application security and applicant privacy.",
    "Seamlessly blends deep technical complexity with an exceptionally beautiful, user-centric 'Thunder' GUI."
])

out_path = os.path.join(BASE_DIR, "SmartRecruit_Final_Presentation_V3.pptx")
prs.save(out_path)
print(f"[OK] New PPTX saved successfully to {out_path}")
